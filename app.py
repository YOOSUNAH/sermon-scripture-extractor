import re, os, io, uuid, time, subprocess, shutil, tempfile, copy
from pathlib import Path
from flask import Flask, request, send_file, render_template, jsonify
from docx import Document
from docx.shared import RGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.oxml.ns import qn as pptx_qn
from lxml import etree
from ppt_generator import (generate_ppt, parse_title_paragraph,
                            extract_verses, split_to_lines, split_to_slides)

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

LIBREOFFICE = (os.environ.get('LIBREOFFICE_PATH')
               or shutil.which('soffice')
               or '/Applications/LibreOffice.app/Contents/MacOS/soffice')
PPT_TEMPLATE = os.path.join(os.path.dirname(__file__), 'ppt_template.pptx')

# 세션 임시 저장소
SESSIONS: dict[str, dict] = {}
SESSION_TTL = 3600  # 1시간


# ── 패턴 ──────────────────────────────────────────────────────────────
RED_PATTERN    = re.compile(r'^\s*\d+:\d+')
YELLOW_PATTERN = re.compile(r'^\s*[가-힣A-Za-z]+\s+\d+:\d+(?![\d-])')


# ── DOCX 처리 ─────────────────────────────────────────────────────────

def _get_or_add_rPr(run):
    rPr = run._r.find(qn('w:rPr'))
    if rPr is None:
        rPr = OxmlElement('w:rPr')
        run._r.insert(0, rPr)
    return rPr


def set_paragraph_red(para):
    for run in para.runs:
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)


def set_paragraph_yellow_highlight(para):
    for run in para.runs:
        rPr = _get_or_add_rPr(run)
        for ex in rPr.findall(qn('w:highlight')):
            rPr.remove(ex)
        hl = OxmlElement('w:highlight')
        hl.set(qn('w:val'), 'yellow')
        rPr.append(hl)


def extract_reference(text: str, is_yellow: bool) -> str | None:
    refs = re.findall(r'(\d+):(\d+)', text)
    if not refs:
        return None
    first_ch, _ = refs[0]
    same_ch_verses = sorted(int(v) for c, v in refs if c == first_ch)
    if len(same_ch_verses) == 1:
        ref = f'{first_ch}:{same_ch_verses[0]}'
    else:
        ref = f'{first_ch}:{same_ch_verses[0]}-{same_ch_verses[-1]}'
    if is_yellow:
        m = re.match(r'^\s*([가-힣A-Za-z]+)\s+', text)
        book = m.group(1) if m else ''
        return f'{book} {ref}'.strip()
    return ref


def process_document(input_bytes: bytes):
    doc = Document(io.BytesIO(input_bytes))
    red_refs, yellow_refs, yellow_para_texts = [], [], []
    title_text = ''

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        if i == 0 and not title_text:
            title_text = text
        if RED_PATTERN.match(text):
            set_paragraph_red(para)
            ref = extract_reference(text, is_yellow=False)
            if ref:
                red_refs.append(ref)
        elif YELLOW_PATTERN.match(text):
            set_paragraph_yellow_highlight(para)
            ref = extract_reference(text, is_yellow=True)
            if ref:
                yellow_refs.append(ref)
            yellow_para_texts.append(text)

    return doc, red_refs, yellow_refs, yellow_para_texts, title_text


def build_refs_text(title_text: str, red_refs: list, yellow_refs: list) -> str:
    passage, sermon_title = parse_title_paragraph(title_text)
    heading = f'{passage} "{sermon_title}"' if sermon_title else title_text

    lines = [heading, '']
    lines += ['본문'] + red_refs + ['']
    lines += ['보조본문'] + yellow_refs
    return '\n'.join(lines)


# ── PDF 변환 ──────────────────────────────────────────────────────────

def docx_to_pdf(docx_path: str, out_dir: str) -> str | None:
    # 1) LibreOffice (headless, 권한 팝업 없음)
    if os.path.exists(LIBREOFFICE):
        try:
            subprocess.run(
                [LIBREOFFICE, '--headless', '--convert-to', 'pdf',
                 '--outdir', out_dir, docx_path],
                capture_output=True, timeout=90
            )
            pdf_path = os.path.join(out_dir, Path(docx_path).stem + '.pdf')
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception:
            pass

    # 2) docx2pdf 폴백 (Word 사용, 권한 팝업 발생 가능)
    try:
        from docx2pdf import convert
        pdf_path = os.path.join(out_dir, Path(docx_path).stem + '.pdf')
        convert(docx_path, pdf_path)
        if os.path.exists(pdf_path):
            return pdf_path
    except Exception:
        pass
    return None


# ── PPT 병합 ──────────────────────────────────────────────────────────

_P_NS   = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'


def _find_section(prs: Presentation, name: str):
    """이름으로 섹션 XML element 반환 (없으면 None)."""
    prs_el = prs.slides._sldIdLst.getparent()
    for sec in prs_el.iter(f'{{{_P14_NS}}}section'):
        if sec.get('name') == name:
            return sec
    return None


def _create_section_after(prs: Presentation, new_name: str, after_name: str):
    """`after_name` 섹션 바로 뒤에 새 섹션을 생성해 반환. after_name이 없으면 None."""
    after_sec = _find_section(prs, after_name)
    if after_sec is None:
        return None
    section_lst = after_sec.getparent()
    new_sec = etree.Element(f'{{{_P14_NS}}}section')
    new_sec.set('name', new_name)
    new_sec.set('id', '{' + str(uuid.uuid4()).upper() + '}')
    etree.SubElement(new_sec, f'{{{_P14_NS}}}sldIdLst')
    after_idx = list(section_lst).index(after_sec)
    section_lst.insert(after_idx + 1, new_sec)
    return new_sec


def _section_last_global_pos(prs: Presentation, section_name: str) -> int:
    """
    섹션의 마지막 슬라이드가 global sldIdLst에서 몇 번째 인덱스인지 반환.
    섹션이 없거나 비어있으면 -1.
    """
    sec = _find_section(prs, section_name)
    if sec is None:
        return -1
    sec_ids = sec.findall(f'{{{_P14_NS}}}sldIdLst/{{{_P14_NS}}}sldId')
    if not sec_ids:
        return -1
    last_sec_id = sec_ids[-1].get('id')

    sldIdLst = prs.slides._sldIdLst
    for i, sldId in enumerate(list(sldIdLst)):
        if sldId.get('id') == last_sec_id:
            return i
    return -1


def merge_ppt(input_ppt_bytes: bytes, generated_ppt_bytes: bytes) -> bytes:
    """
    input_ppt의 '보조본문' 섹션에 generated_ppt 슬라이드를 삽입.
    '보조본문' 섹션이 없고 '본문' 섹션이 있으면 '본문' 뒤에 '보조본문' 섹션을 새로 만들어 추가.
    """
    base_prs = Presentation(io.BytesIO(input_ppt_bytes))
    gen_prs  = Presentation(io.BytesIO(generated_ppt_bytes))

    sldIdLst = base_prs.slides._sldIdLst

    # 삽입 위치 결정: '본문' 섹션 끝 바로 다음
    insert_pos = _section_last_global_pos(base_prs, '본문')
    if insert_pos == -1:
        insert_pos = len(base_prs.slides) - 1
    insert_pos += 1   # 본문 마지막 슬라이드 다음

    # '보조본문' 섹션이 없으면 '본문' 다음에 새로 생성
    bojobon_sec = _find_section(base_prs, '보조본문')
    if bojobon_sec is None:
        bojobon_sec = _create_section_after(base_prs, '보조본문', '본문')

    bojobon_sec_sldIdLst = None
    if bojobon_sec is not None:
        bojobon_sec_sldIdLst = bojobon_sec.find(f'{{{_P14_NS}}}sldIdLst')
        if bojobon_sec_sldIdLst is None:
            bojobon_sec_sldIdLst = etree.SubElement(
                bojobon_sec, f'{{{_P14_NS}}}sldIdLst')

    for gen_slide in gen_prs.slides:
        layout = base_prs.slide_layouts[0]
        new_slide = base_prs.slides.add_slide(layout)

        # 레이아웃 자동 도형 제거
        sp_tree = new_slide.shapes._spTree
        for child in list(sp_tree):
            if child.tag not in (pptx_qn('p:nvGrpSpPr'), pptx_qn('p:grpSpPr')):
                sp_tree.remove(child)

        # 생성 슬라이드 도형 복사
        for shape in gen_slide.shapes:
            sp_tree.append(copy.deepcopy(shape.element))

        # 전환 효과 제거
        for t in new_slide._element.findall(f'{{{_P_NS}}}transition'):
            new_slide._element.remove(t)

        # 방금 추가된 sldId 요소 가져오기
        added_sldId = sldIdLst[-1]
        new_id = added_sldId.get('id')

        # global sldIdLst에서 올바른 위치로 이동
        sldIdLst.remove(added_sldId)
        sldIdLst.insert(insert_pos, added_sldId)
        insert_pos += 1

        # '보조본문' 섹션 XML에 슬라이드 ID 등록
        if bojobon_sec_sldIdLst is not None:
            sec_sldId = etree.SubElement(
                bojobon_sec_sldIdLst, f'{{{_P14_NS}}}sldId')
            sec_sldId.set('id', new_id)

    buf = io.BytesIO()
    base_prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── 세션 정리 ──────────────────────────────────────────────────────────

def _cleanup_old_sessions():
    now = time.time()
    for sid in list(SESSIONS.keys()):
        if now - SESSIONS[sid]['created'] > SESSION_TTL:
            d = SESSIONS[sid].get('dir')
            if d and os.path.isdir(d):
                shutil.rmtree(d, ignore_errors=True)
            del SESSIONS[sid]


# ── 라우트 ────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/process', methods=['POST'])
def process():
    _cleanup_old_sessions()

    if 'file' not in request.files:
        return jsonify({'error': '파일이 없습니다.'}), 400
    f = request.files['file']
    if not f.filename.lower().endswith('.docx'):
        return jsonify({'error': 'DOCX 파일만 업로드 가능합니다.'}), 400

    base_name = Path(f.filename).stem   # 예: 2801_설교문_input
    # 커스텀 파일명 (빈 문자열이면 기본값 사용)
    custom_name = request.form.get('custom_name', '').strip()
    out_stem = custom_name if custom_name else base_name

    input_bytes = f.read()

    doc, red_refs, yellow_refs, yellow_para_texts, title_text = process_document(input_bytes)

    tmp_dir = tempfile.mkdtemp()
    session_id = str(uuid.uuid4())

    # 1) DOCX 저장
    docx_out_name = f'{out_stem}.docx'
    docx_path = os.path.join(tmp_dir, docx_out_name)
    doc.save(docx_path)

    # 2) PDF
    pdf_name = f'{out_stem}.pdf'
    pdf_path = docx_to_pdf(docx_path, tmp_dir)
    if pdf_path:
        final_pdf = os.path.join(tmp_dir, pdf_name)
        if pdf_path != final_pdf:
            os.rename(pdf_path, final_pdf)
        pdf_path = final_pdf

    # 3) 보조본문 PPT 생성
    passage, sermon_title = parse_title_paragraph(title_text)
    ppt_name = f'{out_stem}_보조본문.pptx'
    ppt_path = None
    generated_ppt_bytes = None
    try:
        generated_ppt_bytes = generate_ppt(PPT_TEMPLATE, passage, sermon_title, yellow_para_texts)
        ppt_path = os.path.join(tmp_dir, ppt_name)
        with open(ppt_path, 'wb') as pf:
            pf.write(generated_ppt_bytes)
    except Exception as e:
        print(f'PPT 생성 오류: {e}')

    # 4) 텍스트
    refs_text = build_refs_text(title_text, red_refs, yellow_refs)

    # 세션 저장
    SESSIONS[session_id] = {
        'created': time.time(),
        'dir': tmp_dir,
        'docx':   docx_path if os.path.exists(docx_path) else None,
        'pdf':    pdf_path  if pdf_path and os.path.exists(pdf_path) else None,
        'ppt':    ppt_path  if ppt_path and os.path.exists(ppt_path) else None,
        'merged': None,
        'docx_name':   docx_out_name,
        'pdf_name':    pdf_name,
        'ppt_name':    ppt_name,
        'merged_name': None,
    }

    return jsonify({
        'session_id':   session_id,
        'refs_text':    refs_text,
        'has_pdf':      bool(SESSIONS[session_id]['pdf']),
        'has_ppt':      bool(SESSIONS[session_id]['ppt']),
        'docx_name':    docx_out_name,
        'pdf_name':     pdf_name,
        'ppt_name':     ppt_name,
    })


@app.route('/merge/<session_id>', methods=['POST'])
def merge(session_id):
    sess = SESSIONS.get(session_id)
    if not sess:
        return jsonify({'error': '세션이 만료됐습니다. 다시 처리해 주세요.'}), 404
    if not sess.get('ppt') or not os.path.exists(sess['ppt']):
        return jsonify({'error': '보조본문 PPT가 없어 병합할 수 없습니다.'}), 400

    input_ppt_file = request.files.get('input_ppt')
    if not input_ppt_file or not input_ppt_file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'PPTX 파일을 업로드해 주세요.'}), 400

    with open(sess['ppt'], 'rb') as f:
        generated_ppt_bytes = f.read()

    try:
        input_ppt_bytes = input_ppt_file.read()
        merged_bytes = merge_ppt(input_ppt_bytes, generated_ppt_bytes)
    except Exception as e:
        return jsonify({'error': f'PPT 병합 오류: {e}'}), 500

    custom_name = request.form.get('custom_name', '').strip()
    input_ppt_stem = Path(input_ppt_file.filename).stem
    merged_stem = custom_name if custom_name else f'최종_{input_ppt_stem}'
    merged_ppt_name = f'{merged_stem}.pptx'
    merged_ppt_path = os.path.join(sess['dir'], merged_ppt_name)
    with open(merged_ppt_path, 'wb') as mf:
        mf.write(merged_bytes)

    sess['merged'] = merged_ppt_path
    sess['merged_name'] = merged_ppt_name

    return jsonify({
        'has_merged':  True,
        'merged_name': merged_ppt_name,
    })


@app.route('/download/<session_id>/<file_type>')
def download(session_id, file_type):
    sess = SESSIONS.get(session_id)
    if not sess:
        return '세션이 만료됐습니다.', 404

    PPTX_MIME = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    file_map = {
        'docx':   (sess['docx'],   sess['docx_name'],   'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
        'pdf':    (sess['pdf'],    sess['pdf_name'],    'application/pdf'),
        'ppt':    (sess['ppt'],    sess['ppt_name'],    PPTX_MIME),
        'merged': (sess.get('merged'), sess.get('merged_name'), PPTX_MIME),
    }
    if file_type not in file_map:
        return '잘못된 요청', 400

    path, name, mime = file_map[file_type]
    if not path or not os.path.exists(path):
        return '파일 없음', 404

    return send_file(path, mimetype=mime, as_attachment=True, download_name=name)


if __name__ == '__main__':
    import sys, signal
    port = int(sys.argv[1]) if len(sys.argv) > 1 else 8000

    # 기존 프로세스가 포트를 점유 중이면 종료
    try:
        result = subprocess.run(
            ['lsof', '-ti', f':{port}'],
            capture_output=True, text=True
        )
        pids = result.stdout.strip().split()
        for pid in pids:
            if pid and int(pid) != os.getpid():
                os.kill(int(pid), signal.SIGKILL)
                print(f'포트 {port} 사용 중인 프로세스 {pid} 종료')
    except Exception:
        pass

    print(f'서버 시작 → http://0.0.0.0:{port}/')
    app.run(host='0.0.0.0', port=port, debug=False)
