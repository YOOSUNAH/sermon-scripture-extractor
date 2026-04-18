"""
PPT 생성기 — 보조본문 슬라이드 + 설교제목 슬라이드
"""
import re
import copy
import io
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree as _etree


def _make_el(tag: str):
    """Clark notation 또는 prefix:local 태그로 lxml Element 생성."""
    return _etree.Element(qn(tag) if ':' in tag else tag)

# ── 폰트/색상 상수 ────────────────────────────────────────────────────
FONT_NAME = '나눔고딕 ExtraBold'
BLUE  = RGBColor(0x00, 0xB0, 0xF0)
GREEN = RGBColor(0x00, 0xB0, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── 줄 분리 설정 ──────────────────────────────────────────────────────
# 슬라이드당 최대 2줄 보장이 최우선 조건.
# CHARS_PER_LINE을 올리면 실제 PPT 박스에서 자동 줄바꿈이 일어나 3줄로 표시되므로,
# 실측 안전값 22로 유지한다.
CHARS_PER_LINE = 22   # 한 줄 최대 글자 수 (한국어 기준, 안전값)
LINES_PER_SLIDE = 2   # 슬라이드당 최대 줄 수 (초과 금지)
MAX_LINE_DISPLAY = 22 # PPT 텍스트박스 한 줄 최대 표시 글자 수 (안전값)
# ref(32pt)가 축소된 덕에 content 40pt 기준 22자 이상이어도 시각적으로 들어감.
# 실측상 실제 텍스트박스 가용폭은 26자 정도 → 인라인 판정은 이 값까지 허용.
INLINE_WITH_REF_MAX = 26
REF_FONT_SIZE = 32    # group_ref(예: "(마 17:22-23)") 폰트 크기
REF_MARKER = '\u0001' # 런 분리용 내부 마커 (출력 안 됨)

# ── 템플릿 슬라이드 인덱스 (ppt_template.pptx 기준) ──────────────────
IDX_VERSE_NORMAL = 0   # 보조본문 일반 (하단 참조 있음)
IDX_VERSE_SLASH  = 1   # 보조본문 마지막 (// 있음)
IDX_TITLE        = 2   # 설교제목


# ── 텍스트 분리 유틸 ──────────────────────────────────────────────────

def split_to_lines(text: str, max_chars: int = CHARS_PER_LINE) -> list[str]:
    """단어 경계를 유지하며 텍스트를 줄로 분리."""
    words = text.split()
    lines, cur, cur_len = [], [], 0
    for word in words:
        wl = len(word)
        test = cur_len + (1 if cur else 0) + wl
        if cur and test > max_chars:
            lines.append(' '.join(cur))
            cur, cur_len = [word], wl
        else:
            cur.append(word)
            cur_len = test
    if cur:
        lines.append(' '.join(cur))
    return lines or ['']


def split_to_slides(lines: list[str], max_lines: int = LINES_PER_SLIDE) -> list[list[str]]:
    """줄 목록을 슬라이드 단위로 그룹화."""
    if not lines:
        return [['']]
    return [lines[i:i+max_lines] for i in range(0, len(lines), max_lines)]


# ── DOCX 파싱 ─────────────────────────────────────────────────────────

def parse_title_paragraph(text: str) -> tuple[str, str]:
    """
    '마 28:1-20 "가서 제자 삼으라"...' → (passage='마 28:1-20', title='가서 제자 삼으라')
    """
    passage_m = re.match(r'^([가-힣A-Za-z]+\s+\d+:\d+(?:-\d+)?)', text.strip())
    passage = passage_m.group(1).strip() if passage_m else ''

    # 한국어/유니코드 따옴표 모두 처리
    title_m = re.search(r'["\u201c\u201d\u0022\uff02](.+?)["\u201c\u201d\u0022\uff02]', text)
    if not title_m:
        title_m = re.search(r'[""](.+?)[""]', text)
    title = title_m.group(1).strip() if title_m else ''
    return passage, title


def extract_verses(para_text: str) -> tuple[str, str, list[tuple[str, str]]]:
    """
    '마 27:51 이에 성소... 27:52 무덤들이...'
    → (book='마', chapter='27', [(verse_num, verse_text), ...])
    """
    book_m = re.match(r'^([가-힣A-Za-z]+)\s+', para_text)
    book = book_m.group(1) if book_m else ''

    tokens = re.split(r'(\d+:\d+)\s*', para_text)
    chapter, verses = None, []
    i = 0
    while i < len(tokens):
        tok = tokens[i].strip()
        if re.match(r'^\d+:\d+$', tok):
            ch, v = tok.split(':')
            if chapter is None:
                chapter = ch
            verse_text = tokens[i+1].strip() if i+1 < len(tokens) else ''
            verses.append((v, verse_text))
            i += 2
        else:
            i += 1

    return book, chapter or '', verses


# ── PPT 슬라이드 생성 ─────────────────────────────────────────────────

_P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _add_slide_from_template(out_prs: Presentation, template_slide) -> object:
    """템플릿 슬라이드의 도형들을 복사해 새 슬라이드를 추가."""
    layout = out_prs.slide_layouts[0]
    new_slide = out_prs.slides.add_slide(layout)
    sp_tree = new_slide.shapes._spTree

    # 레이아웃에서 자동으로 추가된 도형 제거
    for child in list(sp_tree):
        if child.tag not in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
            sp_tree.remove(child)

    # 템플릿 도형 복사
    for shape in template_slide.shapes:
        sp_tree.append(copy.deepcopy(shape.element))

    # 전환 효과 제거
    for t in new_slide._element.findall(f'{{{_P_NS}}}transition'):
        new_slide._element.remove(t)

    return new_slide


def _replace_tf_text(tf, text: str):
    """
    텍스트 프레임의 텍스트만 교체 — 포맷팅(색상·폰트·lang 등) 완전 보존.
    템플릿에서 복사된 run의 XML 구조를 그대로 유지한 채 텍스트 내용만 바꾼다.
    """
    # 첫 번째 paragraph의 모든 run 텍스트를 합쳐서 첫 run 하나로 만듦
    paras = tf.paragraphs
    if not paras:
        return
    first_para = paras[0]
    runs = first_para.runs
    if runs:
        # 첫 번째 run에 전체 텍스트 설정
        runs[0]._r.find(qn('a:t')).text = text
        # 나머지 run 제거
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        # run이 없으면 a:t 직접 추가
        r_el = _make_el('a:r')
        t_el = _make_el('a:t')
        t_el.text = text
        r_el.append(t_el)
        first_para._p.append(r_el)
    # 두 번째 이후 paragraph 제거
    for para in paras[1:]:
        para._p.getparent().remove(para._p)


def _get_shape(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def _set_tf(tf, lines: list[str], font_size: int = None, color: RGBColor = None,
            ref_size: int = None):
    """
    텍스트 프레임의 내용을 교체.
    line에 REF_MARKER가 있으면 그 이후 문자열을 ref_size 크기로 별도 run 렌더링.
    """
    def _style(run, size):
        run.font.name = FONT_NAME
        if size:
            run.font.size = Pt(size)
        if color:
            run.font.color.rgb = color

    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if REF_MARKER in line:
            before, after = line.split(REF_MARKER, 1)
            if before:
                r1 = para.add_run()
                r1.text = before
                _style(r1, font_size)
            r2 = para.add_run()
            r2.text = after
            _style(r2, ref_size if ref_size else font_size)
        else:
            run = para.add_run()
            run.text = line
            _style(run, font_size)


def _make_verse_slide(out_prs, template_normal, template_slash,
                      verse_num: str, content_lines: list[str],
                      bottom_text: str | None, use_slash_pos: bool):
    """
    보조본문 슬라이드 생성.
    bottom_text=None → 하단 참조 없음 (중간 연속 슬라이드)
    use_slash_pos=True → // 높이 위치 (살짝 위)
    """
    tmpl = template_slash if use_slash_pos and bottom_text else template_normal
    slide = _add_slide_from_template(out_prs, tmpl)

    # 절 번호 (TextBox 7)
    tb7 = _get_shape(slide, 'TextBox 7')
    if tb7 and tb7.has_text_frame:
        _set_tf(tb7.text_frame, [verse_num], font_size=36, color=WHITE)

    # 절 내용 (TextBox 6) — group_ref는 REF_FONT_SIZE로 축소
    tb6 = _get_shape(slide, 'TextBox 6')
    if tb6 and tb6.has_text_frame:
        _set_tf(tb6.text_frame, content_lines, font_size=40, color=WHITE,
                ref_size=REF_FONT_SIZE)

    # 하단 참조 (TextBox 2)
    tb2 = _get_shape(slide, 'TextBox 2')
    if bottom_text is None:
        # 중간 연속 슬라이드: 하단 텍스트박스 제거
        if tb2:
            tb2.element.getparent().remove(tb2.element)
    else:
        if tb2 and tb2.has_text_frame:
            _set_tf(tb2.text_frame, [bottom_text], font_size=80, color=BLUE)

    return slide


def _make_title_slide(out_prs, template_title, passage: str, sermon_title: str):
    """설교제목 슬라이드 생성."""
    slide = _add_slide_from_template(out_prs, template_title)

    # 텍스트만 교체 — 템플릿의 schemeClr bg1, lang, a:latin, a:ea 등 포맷 그대로 유지
    tb3 = _get_shape(slide, 'TextBox 3')
    if tb3 and tb3.has_text_frame:
        _replace_tf_text(tb3.text_frame, sermon_title)

    tb4 = _get_shape(slide, 'TextBox 4')
    if tb4 and tb4.has_text_frame:
        _replace_tf_text(tb4.text_frame, passage)

    # TextBox 1 ("설교제목", 초록색)는 템플릿에서 그대로 유지
    return slide


# ── 메인 생성 함수 ────────────────────────────────────────────────────

def generate_ppt(template_path: str,
                 passage: str,
                 sermon_title: str,
                 yellow_para_texts: list[str]) -> bytes:
    """
    PPT 바이트 생성.

    template_path: ppt_template.pptx 경로
    passage: '마 28:1-20'
    sermon_title: '가서 제자 삼으라'
    yellow_para_texts: 보조본문 단락 텍스트 목록 (순서대로)
    """
    template_prs = Presentation(template_path)
    t_normal = template_prs.slides[IDX_VERSE_NORMAL]
    t_slash  = template_prs.slides[IDX_VERSE_SLASH]
    t_title  = template_prs.slides[IDX_TITLE]

    # 출력 프레젠테이션 (템플릿 기반으로 마스터/테마 유지)
    out_prs = Presentation(template_path)
    # 기존 슬라이드 제거
    slide_id_lst = out_prs.slides._sldIdLst
    for sldId in list(slide_id_lst):
        rId = sldId.get(qn('r:id'))
        slide_id_lst.remove(sldId)
        try:
            out_prs.slides.part.drop_rel(rId)
        except Exception:
            pass

    # ── 슬라이드 빌드 ──────────────────────────────────────────────
    # 구조: [GROUP1 SLIDES] [설교제목] [GROUP2 SLIDES] [설교제목] ...

    for group_idx, para_text in enumerate(yellow_para_texts):
        book, chapter, verses = extract_verses(para_text)
        if not verses:
            continue

        # 그룹 전체 범위 참조 (예: "마 27:51-52")
        all_verse_nums = [int(v) for v, _ in verses]
        v_min, v_max = min(all_verse_nums), max(all_verse_nums)
        if v_min == v_max:
            group_ref = f'({book} {chapter}:{v_min})'
        else:
            group_ref = f'({book} {chapter}:{v_min}-{v_max})'

        # 이 그룹의 모든 슬라이드 정보를 계산 (carry-over 적용)
        # 마지막 chunk가 LINES_PER_SLIDE 미만이면 다음 절에 단어를 넘겨 재결합 → 빈 슬라이드 방지
        group_slides = []
        carry_words = []  # 이전 절의 마지막 부분 단어들

        for vi, (verse_num, verse_text) in enumerate(verses):
            is_last_verse = (vi == len(verses) - 1)

            # carry 단어 + 현재 절 단어 결합
            combined_words = carry_words + verse_text.split()
            carry_words = []

            lines = split_to_lines(' '.join(combined_words))
            chunks = split_to_slides(lines)

            # 마지막 절이 아니고, 마지막 chunk가 한 줄 미만이면 다음 절로 carry
            if not is_last_verse and chunks and len(chunks[-1]) < LINES_PER_SLIDE:
                carry_words = ' '.join(chunks[-1]).split()
                chunks = chunks[:-1]
                if not chunks:
                    continue  # 이 절 전체를 다음 절로 넘김

            # 마지막 절의 마지막 chunk에 group_ref 추가 (REF_MARKER로 런 분리 → 축소 렌더)
            if is_last_verse and chunks:
                # 괄호 안 공백을 non-breaking space로 치환 → PPT 내부에서 줄 분리 방지
                group_ref_safe = group_ref.replace(' ', '\u00A0')
                last_line = chunks[-1][-1]

                # group_ref는 REF_FONT_SIZE(32pt) 축소 + 괄호/숫자 위주 narrow chars.
                # 실측상 content(40pt 한글)의 약 50% 시각 폭을 차지 → 이를 반영해 인라인 판정.
                # 한도는 INLINE_WITH_REF_MAX(실측 가용폭)까지 허용 → 2줄 이내로 최대한 수렴.
                ref_visual_width = len(group_ref_safe) * 0.5

                if len(last_line) + 1 + ref_visual_width <= INLINE_WITH_REF_MAX:
                    # group_ref 전체가 같은 줄에 들어감 → 인라인 유지
                    chunks[-1][-1] = last_line + ' ' + REF_MARKER + group_ref_safe
                elif len(chunks[-1]) == 1:
                    # 내용 1줄 + 인라인 불가 → group_ref를 2번째 줄로 (잘림 방지)
                    chunks[-1].append(REF_MARKER + group_ref_safe)
                else:
                    # 내용 2줄 + 인라인 불가 → 마지막 줄을 새 chunk로 옮기고 ref와 함께
                    moved = chunks[-1].pop()
                    chunks.append([moved, REF_MARKER + group_ref_safe])

            for ci, chunk in enumerate(chunks):
                is_first_chunk = (ci == 0)
                is_last_chunk = (ci == len(chunks) - 1)
                is_last_of_group = is_last_verse and is_last_chunk

                group_slides.append({
                    'verse_num': verse_num,
                    'content': list(chunk),
                    'is_first_chunk': is_first_chunk,
                    'is_last_of_group': is_last_of_group,
                    'book': book,
                    'chapter': chapter,
                })

        # 슬라이드 생성
        for si, sd in enumerate(group_slides):
            verse_num = sd['verse_num']
            content = sd['content']
            is_first = sd['is_first_chunk']
            is_last_of_group = sd['is_last_of_group']
            bk = sd['book']
            ch = sd['chapter']

            if is_last_of_group:
                # 그룹 마지막: // 스타일
                if is_first:
                    bottom_text = f'{bk} {ch}:{verse_num} //'
                else:
                    bottom_text = '//'
                use_slash = True
            elif is_first:
                # 절의 첫 슬라이드: 참조 표시
                bottom_text = f'{bk} {ch}:{verse_num}'
                use_slash = False
            else:
                # 중간 연속 슬라이드: 참조 없음
                bottom_text = None
                use_slash = False

            _make_verse_slide(out_prs, t_normal, t_slash,
                              verse_num, content, bottom_text, use_slash)

        # 그룹 뒤 설교제목 슬라이드
        _make_title_slide(out_prs, t_title, passage, sermon_title)

    # bytes로 반환
    buf = io.BytesIO()
    out_prs.save(buf)
    buf.seek(0)
    return buf.read()
